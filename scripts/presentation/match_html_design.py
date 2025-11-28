#!/usr/bin/env python3
"""
Convert HTML to PowerPoint matching exact design:
- White slide backgrounds
- Dark text
- Purple accents
- Proper layouts (grids, cards, etc.)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re


# Exact colors from HTML
COLORS = {
    'bg_gradient_1': RGBColor(102, 126, 234),   # #667eea
    'bg_gradient_2': RGBColor(118, 75, 162),    # #764ba2
    'slide_bg': RGBColor(255, 255, 255),        # white
    'title': RGBColor(45, 55, 72),              # #2d3748
    'heading': RGBColor(74, 85, 104),           # #4a5568
    'text': RGBColor(74, 85, 104),              # #4a5568
    'subtitle': RGBColor(113, 128, 150),        # #718096
    'purple': RGBColor(102, 126, 234),          # #667eea
    'light_bg': RGBColor(247, 250, 252),        # #f7fafc
    'green': RGBColor(72, 187, 120),            # #48bb78
}


def parse_slides(html_file):
    """Parse HTML and extract slide data"""
    with open(html_file, 'r') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')

    slides_data = []
    for slide_div in soup.find_all('div', class_='slide'):
        data = {
            'title': '',
            'subtitle': '',
            'emoji': '',
            'content': [],
            'lists': [],
            'stats': [],
            'features': [],
            'impact_cards': [],
            'is_title_slide': False
        }

        # Check if title slide
        if slide_div.find(class_='title-slide'):
            data['is_title_slide'] = True

        # Extract emoji
        emoji_div = slide_div.find(class_='emoji')
        if emoji_div:
            data['emoji'] = emoji_div.get_text(strip=True)

        # Title
        h1 = slide_div.find('h1')
        h2 = slide_div.find('h2')
        if h1:
            data['title'] = h1.get_text(strip=True)
        elif h2:
            data['title'] = h2.get_text(strip=True)

        # Subtitle
        subtitle = slide_div.find(class_='subtitle')
        if subtitle:
            data['subtitle'] = subtitle.get_text(strip=True)

        # Stats
        stats_grid = slide_div.find(class_='stats-grid')
        if stats_grid:
            for stat_card in stats_grid.find_all(class_='stat-card'):
                num = stat_card.find(class_='stat-number')
                label = stat_card.find(class_='stat-label')
                if num and label:
                    data['stats'].append({
                        'number': num.get_text(strip=True),
                        'label': label.get_text(strip=True)
                    })

        # Impact cards
        impact_grid = slide_div.find(class_='impact-grid')
        if impact_grid:
            for card in impact_grid.find_all(class_='impact-card'):
                h4 = card.find('h4')
                p = card.find('p')
                if h4 and p:
                    data['impact_cards'].append({
                        'title': h4.get_text(strip=True),
                        'text': p.get_text(strip=True)
                    })

        # Feature list
        feature_list = slide_div.find(class_='feature-list')
        if feature_list:
            for li in feature_list.find_all('li'):
                data['features'].append(li.get_text(strip=True))

        # Regular lists
        for ul in slide_div.find_all('ul'):
            if 'feature-list' not in ul.get('class', []):
                items = [li.get_text(strip=True) for li in ul.find_all('li')]
                if items:
                    data['lists'].append(items)

        # Paragraphs
        for p in slide_div.find_all('p'):
            classes = p.get('class', [])
            if 'subtitle' not in classes:
                text = p.get_text(strip=True)
                if text and len(text) > 15:
                    data['content'].append(text)

        slides_data.append(data)

    return slides_data


def add_purple_gradient_background(slide):
    """Add purple gradient background (for slide background)"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = COLORS['bg_gradient_1']
    fill.gradient_stops[1].color.rgb = COLORS['bg_gradient_2']


def create_slide_with_white_content(prs, data):
    """Create slide with gradient bg and white content box (matching HTML)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_purple_gradient_background(slide)

    # White content box (like HTML)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(6.5)

    # Add white rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['slide_bg']
    shape.line.fill.background()
    shape.shadow.inherit = False

    return slide, left, top, width, height


def create_title_slide(prs, data):
    """Create title slide matching HTML"""
    slide, box_left, box_top, box_width, box_height = create_slide_with_white_content(prs, data)

    # Emoji
    if data['emoji']:
        emoji_box = slide.shapes.add_textbox(
            box_left, Inches(1.5), box_width, Inches(0.8)
        )
        p = emoji_box.text_frame.paragraphs[0]
        p.text = data['emoji']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.font.size = Pt(60)

    # Title
    title_box = slide.shapes.add_textbox(
        box_left + Inches(0.5), Inches(2.0), box_width - Inches(1), Inches(1.4)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = data['title']
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['title']

    # Subtitle
    if data['subtitle']:
        sub_box = slide.shapes.add_textbox(
            box_left + Inches(0.5), Inches(3.5), box_width - Inches(1), Inches(0.8)
        )
        p = sub_box.text_frame.paragraphs[0]
        p.text = data['subtitle']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.font.size = Pt(26)
        p.font.color.rgb = COLORS['subtitle']

    # GitHub Link
    github_box = slide.shapes.add_textbox(
        box_left + Inches(2.5), Inches(4.4), Inches(4), Inches(0.5)
    )
    p = github_box.text_frame.paragraphs[0]
    p.text = "⭐ View on GitHub: github.com/Siavashghaffari/Code_review_assistant"
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['purple']

    # Stats grid (3 columns)
    if data['stats']:
        stat_top = Inches(4.8)
        stat_width = Inches(2.5)
        stat_height = Inches(1.2)
        spacing = 0.3
        start = box_left + (box_width - (3 * stat_width + 2 * spacing)) / 2

        for i, stat in enumerate(data['stats'][:3]):
            x = start + i * (stat_width + spacing)

            # Stat card with gradient
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, stat_top, Inches(stat_width), Inches(stat_height)
            )
            card.fill.gradient()
            card.fill.gradient_angle = 135
            card.fill.gradient_stops[0].color.rgb = COLORS['bg_gradient_1']
            card.fill.gradient_stops[1].color.rgb = COLORS['bg_gradient_2']
            card.line.fill.background()

            tf = card.text_frame
            tf.clear()

            # Number
            p = tf.paragraphs[0]
            p.text = stat['number']
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            # Label
            p = tf.add_paragraph()
            p.text = stat['label']
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)


def create_content_slide(prs, data):
    """Create content slide matching HTML"""
    slide, box_left, box_top, box_width, box_height = create_slide_with_white_content(prs, data)

    current_y = box_top + Inches(0.5)

    # Title
    title_box = slide.shapes.add_textbox(
        box_left + Inches(0.6), current_y, box_width - Inches(1.2), Inches(0.9)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = data['title']
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['heading']
    current_y += Inches(1.1)

    # Emoji
    if data['emoji']:
        emoji_box = slide.shapes.add_textbox(
            box_left + Inches(0.6), current_y, Inches(1), Inches(0.6)
        )
        p = emoji_box.text_frame.paragraphs[0]
        p.text = data['emoji']
        p.font.size = Pt(48)
        current_y += Inches(0.7)

    # Impact cards (2x2 grid)
    if data['impact_cards']:
        card_width = Inches(4)
        card_height = Inches(1.2)
        spacing = 0.4
        start_x = box_left + Inches(0.8)

        for i, card_data in enumerate(data['impact_cards'][:4]):
            row = i // 2
            col = i % 2
            x = start_x + col * (card_width + spacing)
            y = current_y + row * (card_height + spacing)

            # Card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, y, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = COLORS['light_bg']
            card.line.color.rgb = COLORS['purple']
            card.line.width = Pt(3)

            # Card text
            tf = card.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.1)
            tf.clear()

            p = tf.paragraphs[0]
            p.text = card_data['title']
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = COLORS['title']

            p = tf.add_paragraph()
            p.text = card_data['text'][:80]
            p.font.size = Pt(15)
            p.font.color.rgb = COLORS['text']

        current_y += 2 * (card_height + spacing) + 0.2

    # Feature list with checkmarks
    if data['features']:
        list_box = slide.shapes.add_textbox(
            box_left + Inches(0.8), current_y, box_width - Inches(1.6), Inches(3)
        )
        tf = list_box.text_frame
        tf.word_wrap = True
        tf.clear()

        for feature in data['features'][:6]:
            p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
            clean = re.sub(r'[^\w\s\-,.:()%<>/!?&]', '', feature)
            p.text = '✓ ' + clean[:100]
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(10)

    # Regular lists
    if data['lists'] and not data['features']:
        list_box = slide.shapes.add_textbox(
            box_left + Inches(0.8), current_y, box_width - Inches(1.6), Inches(4)
        )
        tf = list_box.text_frame
        tf.word_wrap = True
        tf.clear()

        for items in data['lists'][:2]:
            for item in items[:8]:
                p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
                clean = re.sub(r'[^\w\s\-,.:()%<>/!?&]', '', item)
                p.text = '• ' + clean[:120]
                p.font.size = Pt(17)
                p.font.color.rgb = COLORS['text']
                p.space_after = Pt(8)


def convert_to_pptx(html_file, output_file):
    """Main conversion"""
    print(f"Converting {html_file} to PowerPoint (matching HTML design)...")

    slides_data = parse_slides(html_file)
    print(f"  Found {len(slides_data)} slides")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, data in enumerate(slides_data):
        print(f"  Creating slide {i+1}: {data['title']}")

        if data['is_title_slide'] or i == 0:
            create_title_slide(prs, data)
        else:
            create_content_slide(prs, data)

    prs.save(output_file)
    size_kb = os.path.getsize(output_file) / 1024
    print(f"✓ Created {output_file} ({size_kb:.1f} KB)")
    print("  Design: Purple gradient background + white content boxes")


if __name__ == '__main__':
    import os
    import subprocess
    import sys

    try:
        from bs4 import BeautifulSoup
    except ImportError:
        print("Installing beautifulsoup4...")
        subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'beautifulsoup4'])
        from bs4 import BeautifulSoup

    convert_to_pptx('presentation_slides.html', 'presentation_slides.pptx')
    print("\n✓ Done! PowerPoint matches HTML design.")
