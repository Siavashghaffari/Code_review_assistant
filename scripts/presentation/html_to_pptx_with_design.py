#!/usr/bin/env python3
"""
Convert HTML presentation to PowerPoint with design preserved
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re


# Color scheme from HTML (purple gradient theme)
COLORS = {
    'primary': RGBColor(102, 126, 234),      # #667eea
    'secondary': RGBColor(118, 75, 162),     # #764ba2
    'dark': RGBColor(45, 55, 72),            # #2d3748
    'text': RGBColor(74, 85, 104),           # #4a5568
    'light': RGBColor(247, 250, 252),        # #f7fafc
    'gold': RGBColor(255, 215, 0),           # #ffd700
    'white': RGBColor(255, 255, 255),
}


def parse_slides(html_file):
    """Extract slide content from HTML"""
    with open(html_file, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')

    slides_data = []
    slides = soup.find_all('div', class_='slide')

    for slide in slides:
        data = {
            'title': '',
            'subtitle': '',
            'content': [],
            'lists': [],
            'stats': [],
        }

        # Title
        h1 = slide.find('h1')
        h2 = slide.find('h2')
        if h1:
            data['title'] = h1.get_text(strip=True)
        elif h2:
            data['title'] = h2.get_text(strip=True)

        # Subtitle
        subtitle = slide.find(class_='subtitle')
        if subtitle:
            data['subtitle'] = subtitle.get_text(strip=True)

        # Large text
        large = slide.find(class_='large-text')
        if large:
            data['content'].append(large.get_text(strip=True))

        # Lists
        for ul in slide.find_all('ul'):
            items = [li.get_text(strip=True) for li in ul.find_all('li')]
            if items:
                data['lists'].append(items)

        # Stats
        stats_grid = slide.find(class_='stats-grid')
        if stats_grid:
            for stat in stats_grid.find_all(class_='stat-box'):
                num = stat.find(class_='stat-number')
                label = stat.find(class_='stat-label')
                if num and label:
                    data['stats'].append({
                        'number': num.get_text(strip=True),
                        'label': label.get_text(strip=True)
                    })

        # Paragraphs
        for p in slide.find_all('p'):
            if 'subtitle' not in p.get('class', []) and 'large-text' not in p.get('class', []):
                text = p.get_text(strip=True)
                if text and len(text) > 10:
                    data['content'].append(text)

        slides_data.append(data)

    return slides_data


def add_gradient_background(slide):
    """Add purple gradient background to slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = COLORS['primary']
    fill.gradient_stops[1].color.rgb = COLORS['secondary']


def create_title_slide(prs, data):
    """Create title slide with design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(slide)

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = data['title']
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Subtitle
    if data['subtitle']:
        top = Inches(4.2)
        height = Inches(0.8)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = data['subtitle']
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']

    # Stats if present
    if data['stats']:
        stat_top = Inches(5.2)
        stat_width = Inches(1.8)
        stat_height = Inches(1)
        spacing = 0.2

        total_width = len(data['stats']) * stat_width + (len(data['stats']) - 1) * spacing
        start_left = (10 - total_width) / 2

        for i, stat in enumerate(data['stats'][:4]):
            left = Inches(start_left + i * (stat_width + spacing))
            shape = slide.shapes.add_shape(1, left, Inches(stat_top), Inches(stat_width), Inches(stat_height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
            shape.fill.fore_color.brightness = -0.25
            shape.line.color.rgb = COLORS['white']

            tf = shape.text_frame
            tf.clear()

            # Number
            p = tf.paragraphs[0]
            p.text = stat['number']
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = COLORS['gold']

            # Label
            p = tf.add_paragraph()
            p.text = stat['label']
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['white']


def create_content_slide(prs, data):
    """Create content slide with design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(slide)

    # Title
    left = Inches(0.5)
    top = Inches(0.4)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data['title']
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Content area
    content_top = Inches(1.5)
    content_height = Inches(5)

    content_box = slide.shapes.add_textbox(left, content_top, width, content_height)
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.clear()

    # Subtitle
    if data['subtitle']:
        p = tf.paragraphs[0]
        p.text = data['subtitle']
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.color.rgb = COLORS['white']
        p.space_after = Pt(12)

    # Content paragraphs
    for content in data['content'][:3]:
        p = tf.add_paragraph() if data['subtitle'] or tf.paragraphs[0].text else tf.paragraphs[0]
        clean_text = re.sub(r'[^\w\s\-,.:()%<>/!?&]', '', content)
        p.text = clean_text[:200]
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.space_after = Pt(8)

    # Lists
    for items in data['lists'][:1]:
        for item in items[:6]:
            p = tf.add_paragraph()
            clean_text = re.sub(r'[^\w\s\-,.:()%<>/!?&]', '', item)
            p.text = '• ' + clean_text[:150]
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['white']
            p.level = 0
            p.space_after = Pt(6)


def convert_html_to_pptx(html_file, output_file):
    """Main conversion function"""
    print(f"Converting {html_file} to PowerPoint with design...")

    slides_data = parse_slides(html_file)
    print(f"  Found {len(slides_data)} slides")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for i, data in enumerate(slides_data):
        print(f"  Creating slide {i+1}: {data['title']}")

        if i == 0:
            create_title_slide(prs, data)
        else:
            create_content_slide(prs, data)

    prs.save(output_file)
    print(f"✓ Saved to {output_file}")


if __name__ == '__main__':
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        import subprocess
        print("Installing beautifulsoup4...")
        subprocess.check_call(['pip', 'install', 'beautifulsoup4'])
        from bs4 import BeautifulSoup

    convert_html_to_pptx('presentation_slides.html', 'presentation_slides.pptx')
    print("\n✓ PowerPoint created with purple gradient design!")
