#!/usr/bin/env python3
"""
Convert HTML presentations to PowerPoint PPTX format
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
import re
from typing import List, Dict


def parse_html_slides(html_file: str) -> List[Dict]:
    """Parse HTML file and extract slide content"""
    with open(html_file, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')

    slides_data = []
    slides = soup.find_all('div', class_='slide')

    for slide in slides:
        slide_content = {
            'title': '',
            'subtitle': '',
            'content': [],
            'code_blocks': [],
            'stats': [],
            'features': []
        }

        # Extract title
        h1 = slide.find('h1')
        h2 = slide.find('h2')
        if h1:
            slide_content['title'] = h1.get_text(strip=True)
        elif h2:
            slide_content['title'] = h2.get_text(strip=True)

        # Extract subtitle
        subtitle = slide.find(class_='subtitle')
        if subtitle:
            slide_content['subtitle'] = subtitle.get_text(strip=True)

        # Extract all text content
        for p in slide.find_all('p'):
            if 'subtitle' not in p.get('class', []):
                text = p.get_text(strip=True)
                if text and text not in slide_content['title']:
                    slide_content['content'].append(text)

        # Extract lists
        for ul in slide.find_all('ul'):
            items = [li.get_text(strip=True) for li in ul.find_all('li')]
            if items:
                slide_content['content'].extend(items)

        # Extract code blocks
        for code in slide.find_all(class_='code-block'):
            code_text = code.get_text(strip=True)
            slide_content['code_blocks'].append(code_text)

        # Extract stats
        stats_grid = slide.find(class_='stats-grid')
        if stats_grid:
            stat_boxes = stats_grid.find_all(class_=['stat-box', 'stat-card'])
            for stat in stat_boxes:
                number = stat.find(class_='stat-number')
                label = stat.find(class_='stat-label')
                if number and label:
                    slide_content['stats'].append({
                        'number': number.get_text(strip=True),
                        'label': label.get_text(strip=True)
                    })

        # Extract features
        feature_grid = slide.find(class_=['feature-grid', 'feature-list', 'impact-grid'])
        if feature_grid:
            features = feature_grid.find_all(['li', 'h3', 'h4'])
            for feat in features:
                text = feat.get_text(strip=True)
                if text and len(text) > 3:
                    slide_content['features'].append(text)

        slides_data.append(slide_content)

    return slides_data


def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide


def create_content_slide(prs, slide_data):
    """Create a content slide"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title = slide.shapes.title
    title.text = slide_data['title']

    # Add content
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.clear()

    # Add subtitle if exists
    if slide_data['subtitle']:
        p = text_frame.paragraphs[0]
        p.text = slide_data['subtitle']
        p.font.size = Pt(18)
        p.font.italic = True
        p.space_after = Pt(12)

    # Add stats if exists
    if slide_data['stats']:
        for stat in slide_data['stats'][:4]:  # Limit to 4 stats
            p = text_frame.add_paragraph()
            p.text = f"{stat['number']} - {stat['label']}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.level = 0

    # Add features
    if slide_data['features']:
        for feature in slide_data['features'][:10]:  # Limit to 10 features
            p = text_frame.add_paragraph()
            # Remove emoji and clean text
            clean_feature = re.sub(r'[^\w\s\-,.:()%<>/]', '', feature)
            p.text = clean_feature[:150]  # Truncate long text
            p.font.size = Pt(14)
            p.level = 0 if not clean_feature.startswith('  ') else 1

    # Add regular content
    for content in slide_data['content'][:8]:  # Limit to 8 content items
        if content and content not in slide_data['subtitle']:
            p = text_frame.add_paragraph()
            clean_content = re.sub(r'[^\w\s\-,.:()%<>/!?]', '', content)
            p.text = clean_content[:200]  # Truncate long text
            p.font.size = Pt(14)
            p.level = 0

    # Add code blocks on separate text box
    if slide_data['code_blocks']:
        left = Inches(0.5)
        top = Inches(4.5)
        width = Inches(9)
        height = Inches(2)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for code in slide_data['code_blocks'][:2]:  # Limit to 2 code blocks
            p = text_frame.add_paragraph()
            p.text = code[:300]  # Truncate long code
            p.font.name = 'Courier New'
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0, 0, 0)

    return slide


def convert_html_to_pptx(html_file: str, output_file: str):
    """Main conversion function"""
    print(f"Converting {html_file} to {output_file}...")

    # Parse HTML
    slides_data = parse_html_slides(html_file)
    print(f"Found {len(slides_data)} slides")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Process slides
    for i, slide_data in enumerate(slides_data):
        print(f"  Processing slide {i+1}: {slide_data['title']}")

        if i == 0 and slide_data['subtitle']:
            # First slide as title slide
            create_title_slide(prs, slide_data['title'], slide_data['subtitle'])
        else:
            # Regular content slide
            create_content_slide(prs, slide_data)

    # Save presentation
    prs.save(output_file)
    print(f"✓ Saved to {output_file}")


if __name__ == '__main__':
    # Install BeautifulSoup if not available
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        print("Installing beautifulsoup4...")
        import subprocess
        subprocess.check_call(['pip', 'install', 'beautifulsoup4'])
        from bs4 import BeautifulSoup

    # Convert both HTML files
    conversions = [
        ('presentation_slides.html', 'presentation_slides.pptx'),
        ('presentation.html', 'presentation.pptx')
    ]

    for html_file, pptx_file in conversions:
        try:
            convert_html_to_pptx(html_file, pptx_file)
        except Exception as e:
            print(f"Error converting {html_file}: {e}")

    print("\n✓ Conversion complete!")
